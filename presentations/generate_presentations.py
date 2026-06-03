#!/usr/bin/env python3
"""Generate KM0 Digital presentations in CA, ES, EN with distinct themes."""

from __future__ import annotations

import subprocess
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

BASE = Path(__file__).resolve().parent
LOGO = BASE / "assets" / "logo.png"
OUT = BASE

# Brand palette
NAVY = RGBColor(0x0B, 0x12, 0x20)
SURFACE = RGBColor(0xF5, 0xF5, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ORANGE = RGBColor(0xFF, 0x5F, 0x2E)
MAGENTA = RGBColor(0xE0, 0x40, 0xA0)
PURPLE = RGBColor(0x7B, 0x3F, 0xE4)
BLUE = RGBColor(0x00, 0x7B, 0xFF)
GRAY = RGBColor(0x6B, 0x72, 0x80)
LIGHT_GRAY = RGBColor(0xD1, 0xD5, 0xDB)
DARK_TEXT = RGBColor(0x1F, 0x29, 0x37)


@dataclass
class Theme:
    name: str
    bg: RGBColor
    title_color: RGBColor
    body_color: RGBColor
    accent: RGBColor
    accent2: RGBColor
    subtitle_color: RGBColor
    slide_style: str  # dark | light | split
    footer_color: RGBColor


THEMES = {
    "ca": Theme(
        name="Origen Local",
        bg=NAVY,
        title_color=WHITE,
        body_color=RGBColor(0xE5, 0xE7, 0xEB),
        accent=ORANGE,
        accent2=MAGENTA,
        subtitle_color=ORANGE,
        slide_style="dark",
        footer_color=GRAY,
    ),
    "es": Theme(
        name="Impacto Digital",
        bg=WHITE,
        title_color=NAVY,
        body_color=RGBColor(0x37, 0x41, 0x51),
        accent=MAGENTA,
        accent2=PURPLE,
        subtitle_color=MAGENTA,
        slide_style="light",
        footer_color=GRAY,
    ),
    "en": Theme(
        name="Sovereign Tech",
        bg=SURFACE,
        title_color=NAVY,
        body_color=DARK_TEXT,
        accent=BLUE,
        accent2=PURPLE,
        subtitle_color=BLUE,
        slide_style="split",
        footer_color=GRAY,
    ),
}

CONTENT = {
    "ca": {
        "filename": "KM0-Origen-Local-CA",
        "lang_label": "Català",
        "slides": [
            {
                "type": "title",
                "title": "Kilòmetre 0 Digital",
                "subtitle": "Moviment digital comunitari",
                "tagline": "RECUPERA EL CONTROL DIGITAL",
                "footer": "km0digital.com · cloud.km0digital.com",
            },
            {
                "type": "content",
                "title": "Per què existeix KM0?",
                "bullets": [
                    "Moviment social i tecnològic construït per i per a la comunitat",
                    "Alternativa real a la dependència de Big Tech",
                    "Serveis oberts, propers i sostenibles, no un altre proveïdor més",
                    "Persones reals darrere del projecte, no un call center anònim",
                ],
            },
            {
                "type": "two_col",
                "title": "L'habitual vs. La proposta KM0",
                "left_title": "L'habitual",
                "left": [
                    "Les teves dades com a producte",
                    "Suport anònim i llunyà",
                    "Preus que no reflecteixen el cost real",
                    "Dependència d'uns pocs proveïdors globals",
                ],
                "right_title": "La proposta KM0",
                "right": [
                    "Control i transparència sobre les teves dades",
                    "Persones reals, a prop teu",
                    "Cost honest i sostenible",
                    "Comunitat oberta a qui vulgui participar",
                ],
            },
            {
                "type": "chips",
                "title": "Principis del moviment",
                "chips": ["Obert", "Transparent", "Comunitari", "Local"],
                "body": "Retornar el control digital a persones, famílies, associacions, cooperatives i centres educatius.",
            },
            {
                "type": "content",
                "title": "KM0 Cloud, ja en producció",
                "bullets": [
                    "Emmagatzematge i col·laboració al núvol com un drive propi",
                    "Comparteix fitxers, sincronitza i treballa en equip",
                    "Edició d'Office al navegador amb Collabora Online (DOCX, XLSX, PPTX)",
                    "Accés des de web, Android i iOS",
                    "URL: cloud.km0digital.com",
                ],
                "highlight": "Les teves dades, el teu origen. El moviment creix amb tu.",
            },
            {
                "type": "content",
                "title": "OpenCloud, la plataforma oberta",
                "bullets": [
                    "Plataforma open source per a fitxers, compartició i col·laboració",
                    "Desenvolupada per OpenCloud GmbH (Heinlein Group)",
                    "Soberania digital: sense dependència obligatòria de núvols externs",
                    "Arquitectura de microserveis en Go, sense base de dades SQL",
                    "KM0 desplega OpenCloud 7.0.0 en producció",
                ],
            },
            {
                "type": "content",
                "title": "Arquitectura tècnica",
                "bullets": [
                    "Debian 13 + Docker Compose sobre servidors Hetzner (UE, Alemanya)",
                    "Nginx com a proxy HTTPS únic; contenidors només a loopback",
                    "OpenCloud 7.0.0 · Dex OIDC · Collabora CODE · WOPI",
                    "Autenticació híbrida: Google, Apple i LDAP local (IDM)",
                    "Certificats TLS Let's Encrypt · UFW + Fail2ban",
                ],
            },
            {
                "type": "content",
                "title": "Autenticació unificada amb Dex",
                "bullets": [
                    "Dex com a emissor OIDC central per a web, escriptori i mòbil",
                    "Login amb Google, Apple o credencials locals d'OpenCloud",
                    "Interfície de login personalitzada KM0 (CA | ES | EN | DE)",
                    "Clients natius: OpenCloud Desktop, Android i iOS",
                ],
            },
            {
                "type": "content",
                "title": "Confiança i seguretat",
                "bullets": [
                    "Infraestructura a la UE (Falkenstein, Alemanya)",
                    "Operat per AMVARA CONSULTING S.L., certificat ISO/IEC 27001:2022",
                    "Trànsit xifrat TLS · sense vendre dades per a publicitat",
                    "Blog tècnic públic: transparència des del dia zero",
                    "Backups documentats i runbook operatiu",
                ],
            },
            {
                "type": "timeline",
                "title": "Diari del projecte (Dies 0–5)",
                "items": [
                    ("Dia 0", "Fonaments del servidor: Debian, Docker, Nginx"),
                    ("Dia 1", "Primer desplegament OpenCloud + web KM0"),
                    ("Dia 2", "Upgrade a OpenCloud 7 + Dex OIDC + backups"),
                    ("Dia 3", "Domini km0digital.com + Collabora + FAQ multilingüe"),
                    ("Dia 4", "Login LDAP local unificat amb IDM d'OpenCloud"),
                    ("Dia 5", "Trobada de visió a Masnou: per què i per a qui"),
                ],
            },
            {
                "type": "content",
                "title": "Competir sense ser Big Tech",
                "bullets": [
                    "Equips petits amb eines modernes (incloent IA) poden competir amb grans estructures",
                    "Preus honestos basats en cost real, no en el producte de les teves dades",
                    "Suport humà i proximitat: ens trobem en persona si cal",
                    "Local = geogràfic (Barcelona/Masnou), operatiu i ja funcionant",
                ],
            },
            {
                "type": "cta",
                "title": "Forma part del moviment",
                "bullets": [
                    "Explora: km0digital.com",
                    "Prova KM0 Cloud: cloud.km0digital.com",
                    "Tutorials: web, Android i iOS",
                    "Contacte: hello.yoel@amvara.de",
                    "Uneix-te al grup de WhatsApp de la comunitat",
                ],
            },
        ],
    },
    "es": {
        "filename": "KM0-Impacto-Digital-ES",
        "lang_label": "Español",
        "slides": [
            {
                "type": "title",
                "title": "Kilómetro 0 Digital",
                "subtitle": "Movimiento digital comunitario",
                "tagline": "RECUPERA EL CONTROL DIGITAL",
                "footer": "km0digital.com · cloud.km0digital.com",
            },
            {
                "type": "content",
                "title": "¿Por qué existe KM0?",
                "bullets": [
                    "Movimiento social y tecnológico construido por y para la comunidad",
                    "Alternativa real a la dependencia de Big Tech",
                    "Servicios abiertos, cercanos y sostenibles, no otro proveedor más",
                    "Personas reales detrás del proyecto, no un call center anónimo",
                ],
            },
            {
                "type": "two_col",
                "title": "Lo habitual vs. La propuesta KM0",
                "left_title": "Lo habitual",
                "left": [
                    "Tus datos como producto",
                    "Soporte anónimo y lejano",
                    "Precios que no reflejan el coste real",
                    "Dependencia de unos pocos proveedores globales",
                ],
                "right_title": "La propuesta KM0",
                "right": [
                    "Control y transparencia sobre tus datos",
                    "Personas reales, cerca de ti",
                    "Coste honesto y sostenible",
                    "Comunidad abierta a quien quiera participar",
                ],
            },
            {
                "type": "chips",
                "title": "Principios del movimiento",
                "chips": ["Abierto", "Transparente", "Comunitario", "Local"],
                "body": "Devolver el control digital a personas, familias, asociaciones, cooperativas y centros educativos.",
            },
            {
                "type": "content",
                "title": "KM0 Cloud, ya en producción",
                "bullets": [
                    "Almacenamiento y colaboración en la nube como un drive propio",
                    "Comparte archivos, sincroniza y trabaja en equipo",
                    "Edición de Office en el navegador con Collabora Online (DOCX, XLSX, PPTX)",
                    "Acceso desde web, Android e iOS",
                    "URL: cloud.km0digital.com",
                ],
                "highlight": "Tus datos, tu origen. El movimiento crece contigo.",
            },
            {
                "type": "content",
                "title": "OpenCloud, la plataforma abierta",
                "bullets": [
                    "Plataforma open source para archivos, compartición y colaboración",
                    "Desarrollada por OpenCloud GmbH (Heinlein Group)",
                    "Soberanía digital: sin dependencia obligatoria de nubes externas",
                    "Arquitectura de microservicios en Go, sin base de datos SQL",
                    "KM0 despliega OpenCloud 7.0.0 en producción",
                ],
            },
            {
                "type": "content",
                "title": "Arquitectura técnica",
                "bullets": [
                    "Debian 13 + Docker Compose en servidores Hetzner (UE, Alemania)",
                    "Nginx como proxy HTTPS único; contenedores solo en loopback",
                    "OpenCloud 7.0.0 · Dex OIDC · Collabora CODE · WOPI",
                    "Autenticación híbrida: Google, Apple y LDAP local (IDM)",
                    "Certificados TLS Let's Encrypt · UFW + Fail2ban",
                ],
            },
            {
                "type": "content",
                "title": "Autenticación unificada con Dex",
                "bullets": [
                    "Dex como emisor OIDC central para web, escritorio y móvil",
                    "Login con Google, Apple o credenciales locales de OpenCloud",
                    "Interfaz de login personalizada KM0 (CA | ES | EN | DE)",
                    "Clientes nativos: OpenCloud Desktop, Android e iOS",
                ],
            },
            {
                "type": "content",
                "title": "Confianza y seguridad",
                "bullets": [
                    "Infraestructura en la UE (Falkenstein, Alemania)",
                    "Operado por AMVARA CONSULTING S.L., certificado ISO/IEC 27001:2022",
                    "Tráfico cifrado TLS · sin vender datos para publicidad",
                    "Blog técnico público: transparencia desde el día cero",
                    "Backups documentados y runbook operativo",
                ],
            },
            {
                "type": "timeline",
                "title": "Diario del proyecto (Días 0–5)",
                "items": [
                    ("Día 0", "Cimientos del servidor: Debian, Docker, Nginx"),
                    ("Día 1", "Primer despliegue OpenCloud + web KM0"),
                    ("Día 2", "Upgrade a OpenCloud 7 + Dex OIDC + backups"),
                    ("Día 3", "Dominio km0digital.com + Collabora + FAQ multilingüe"),
                    ("Día 4", "Login LDAP local unificado con IDM de OpenCloud"),
                    ("Día 5", "Encuentro de visión en Masnou: por qué y para quién"),
                ],
            },
            {
                "type": "content",
                "title": "Competir sin ser Big Tech",
                "bullets": [
                    "Equipos pequeños con herramientas modernas (incluida IA) pueden competir con grandes estructuras",
                    "Precios honestos basados en coste real, no en el producto de tus datos",
                    "Soporte humano y proximidad: nos vemos en persona si hace falta",
                    "Local = geográfico (Barcelona/Masnou), operativo y ya funcionando",
                ],
            },
            {
                "type": "cta",
                "title": "Forma parte del movimiento",
                "bullets": [
                    "Explora: km0digital.com",
                    "Prueba KM0 Cloud: cloud.km0digital.com",
                    "Tutoriales: web, Android e iOS",
                    "Contacto: hello.yoel@amvara.de",
                    "Únete al grupo de WhatsApp de la comunidad",
                ],
            },
        ],
    },
    "en": {
        "filename": "KM0-Sovereign-Tech-EN",
        "lang_label": "English",
        "slides": [
            {
                "type": "title",
                "title": "Kilometer 0 Digital",
                "subtitle": "Community digital movement",
                "tagline": "RECLAIM YOUR DIGITAL CONTROL",
                "footer": "km0digital.com · cloud.km0digital.com",
            },
            {
                "type": "content",
                "title": "Why does KM0 exist?",
                "bullets": [
                    "A social and technological movement built by and for the community",
                    "A real alternative to Big Tech dependency",
                    "Open, local, and sustainable services, not just another provider",
                    "Real people behind the project, not an anonymous call centre",
                ],
            },
            {
                "type": "two_col",
                "title": "Business as usual vs. The KM0 proposal",
                "left_title": "The usual way",
                "left": [
                    "Your data as the product",
                    "Anonymous, distant support",
                    "Prices that do not reflect real costs",
                    "Total dependence on a few global providers",
                ],
                "right_title": "The KM0 proposal",
                "right": [
                    "Control and transparency over your data",
                    "Real people, close to you",
                    "Honest, sustainable pricing",
                    "An open community anyone can join",
                ],
            },
            {
                "type": "chips",
                "title": "Movement principles",
                "chips": ["Open", "Transparent", "Community-driven", "Local"],
                "body": "Returning digital control to people, families, associations, cooperatives, and schools.",
            },
            {
                "type": "content",
                "title": "KM0 Cloud, live in production",
                "bullets": [
                    "Cloud storage and collaboration like your own drive",
                    "Share files, sync, and work as a team",
                    "In-browser Office editing with Collabora Online (DOCX, XLSX, PPTX)",
                    "Access from web, Android, and iOS",
                    "URL: cloud.km0digital.com",
                ],
                "highlight": "Your data, your origin. The movement grows with you.",
            },
            {
                "type": "content",
                "title": "OpenCloud, the open platform",
                "bullets": [
                    "Open-source platform for files, sharing, and collaboration",
                    "Built by OpenCloud GmbH (Heinlein Group)",
                    "Digital sovereignty: no mandatory external cloud dependency",
                    "Go microservices architecture, no SQL database",
                    "KM0 runs OpenCloud 7.0.0 in production",
                ],
            },
            {
                "type": "content",
                "title": "Technical architecture",
                "bullets": [
                    "Debian 13 + Docker Compose on Hetzner servers (EU, Germany)",
                    "Nginx as sole HTTPS proxy; containers on loopback only",
                    "OpenCloud 7.0.0 · Dex OIDC · Collabora CODE · WOPI",
                    "Hybrid auth: Google, Apple, and local LDAP (IDM)",
                    "Let's Encrypt TLS · UFW + Fail2ban hardening",
                ],
            },
            {
                "type": "content",
                "title": "Unified authentication with Dex",
                "bullets": [
                    "Dex as central OIDC issuer for web, desktop, and mobile",
                    "Sign in with Google, Apple, or local OpenCloud credentials",
                    "KM0-branded login UI (CA | ES | EN | DE)",
                    "Native clients: OpenCloud Desktop, Android, and iOS",
                ],
            },
            {
                "type": "content",
                "title": "Trust and security",
                "bullets": [
                    "EU infrastructure (Falkenstein, Germany)",
                    "Operated by AMVARA CONSULTING S.L., ISO/IEC 27001:2022 certified",
                    "TLS-encrypted traffic · no selling data for advertising",
                    "Public technical blog: transparency from day zero",
                    "Documented backups and operational runbook",
                ],
            },
            {
                "type": "timeline",
                "title": "Project diary (Days 0–5)",
                "items": [
                    ("Day 0", "Server foundations: Debian, Docker, Nginx"),
                    ("Day 1", "First OpenCloud deployment + KM0 web site"),
                    ("Day 2", "Upgrade to OpenCloud 7 + Dex OIDC + backups"),
                    ("Day 3", "km0digital.com domain + Collabora + multilingual FAQ"),
                    ("Day 4", "Unified local LDAP login with OpenCloud IDM"),
                    ("Day 5", "Vision meeting in Masnou: why and for whom"),
                ],
            },
            {
                "type": "content",
                "title": "Competing without being Big Tech",
                "bullets": [
                    "Small teams with modern tools (including AI) can compete with large structures",
                    "Honest pricing based on real cost, not on monetising your data",
                    "Human support and proximity: meet in person if needed",
                    "Local = geographic (Barcelona/Masnou), operational, and already running",
                ],
            },
            {
                "type": "cta",
                "title": "Join the movement",
                "bullets": [
                    "Explore: km0digital.com",
                    "Try KM0 Cloud: cloud.km0digital.com",
                    "Tutorials: web, Android, and iOS",
                    "Contact: hello.yoel@amvara.de",
                    "Join the community WhatsApp group",
                ],
            },
        ],
    },
}


def set_slide_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_accent_bar(slide, theme: Theme, y: float = 0, height: float = 0.08) -> None:
    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(y),
        Inches(13.333),
        Inches(height),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = theme.accent
    bar.line.fill.background()


def add_split_sidebar(slide, theme: Theme) -> None:
    sidebar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(2.2),
        Inches(7.5),
    )
    sidebar.fill.solid()
    sidebar.fill.fore_color.rgb = NAVY
    sidebar.line.fill.background()

    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(2.2),
        Inches(0),
        Inches(0.06),
        Inches(7.5),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = theme.accent
    accent.line.fill.background()


def add_logo(slide, left: float, top: float, height: float = 1.4) -> None:
    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(left), Inches(top), height=Inches(height))


def add_footer(slide, theme: Theme, text: str, lang: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.4))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{text}  ·  Theme: {theme.name}  ·  {lang}"
    p.font.size = Pt(9)
    p.font.color.rgb = theme.footer_color
    p.alignment = PP_ALIGN.RIGHT


def style_title(text_frame, theme: Theme, size: int = 32, left_margin: float = 0.8) -> None:
    for i, para in enumerate(text_frame.paragraphs):
        para.font.bold = True
        para.font.size = Pt(size if i == 0 else size - 4)
        para.font.color.rgb = theme.title_color
        para.font.name = "Inter"
        para.space_after = Pt(6)


def style_body(text_frame, theme: Theme, size: int = 16) -> None:
    for para in text_frame.paragraphs:
        para.font.size = Pt(size)
        para.font.color.rgb = theme.body_color
        para.font.name = "Inter"
        para.space_after = Pt(8)
        para.level = 0


def add_bullets(text_frame, items: list[str], theme: Theme, size: int = 16) -> None:
    text_frame.clear()
    for i, item in enumerate(items):
        p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = theme.body_color
        p.font.name = "Inter"
        p.space_after = Pt(10)


def content_left(theme: Theme) -> float:
    return 2.8 if theme.slide_style == "split" else 0.8


def build_title_slide(prs: Presentation, slide_data: dict, theme: Theme, lang: str) -> None:
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    set_slide_bg(slide, theme.bg if theme.slide_style == "dark" else NAVY)

    if theme.slide_style == "light":
        add_accent_bar(slide, theme, y=6.8, height=0.7)

    add_logo(slide, 5.5, 0.6, 1.8)

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(11.7), Inches(1.2))
    tf = title_box.text_frame
    tf.text = slide_data["title"]
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.size = Pt(40)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.name = "Inter"

    sub = slide.shapes.add_textbox(Inches(0.8), Inches(3.5), Inches(11.7), Inches(0.5))
    stf = sub.text_frame
    stf.text = slide_data["subtitle"]
    stf.paragraphs[0].alignment = PP_ALIGN.CENTER
    stf.paragraphs[0].font.size = Pt(18)
    stf.paragraphs[0].font.color.rgb = theme.subtitle_color if theme.slide_style != "dark" else ORANGE
    stf.paragraphs[0].font.name = "Inter"

    tag = slide.shapes.add_textbox(Inches(0.8), Inches(4.3), Inches(11.7), Inches(0.8))
    ttf = tag.text_frame
    ttf.text = slide_data["tagline"]
    ttf.paragraphs[0].alignment = PP_ALIGN.CENTER
    ttf.paragraphs[0].font.size = Pt(22)
    ttf.paragraphs[0].font.bold = True
    ttf.paragraphs[0].font.color.rgb = MAGENTA if theme.slide_style == "light" else theme.accent2
    ttf.paragraphs[0].font.name = "Inter"

    add_footer(slide, theme, slide_data["footer"], lang)


def build_content_slide(
    prs: Presentation, slide_data: dict, theme: Theme, lang: str, slide_num: int
) -> None:
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    set_slide_bg(slide, theme.bg)

    left = content_left(theme)
    if theme.slide_style == "light":
        add_accent_bar(slide, theme, y=0, height=0.12)
    elif theme.slide_style == "split":
        add_split_sidebar(slide, theme)
        add_logo(slide, 0.35, 0.3, 1.0)

    title_box = slide.shapes.add_textbox(Inches(left), Inches(0.35), Inches(11.5 - left), Inches(0.9))
    tf = title_box.text_frame
    tf.text = slide_data["title"]
    style_title(tf, theme, size=30)

    body_top = 1.35
    body_box = slide.shapes.add_textbox(Inches(left), Inches(body_top), Inches(11.5 - left), Inches(5.2))
    btf = body_box.text_frame
    btf.word_wrap = True
    add_bullets(btf, slide_data["bullets"], theme)

    if slide_data.get("highlight"):
        hi = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(6.0),
            Inches(11.0 - left),
            Inches(0.7),
        )
        hi.fill.solid()
        hi.fill.fore_color.rgb = theme.accent
        hi.line.fill.background()
        htf = hi.text_frame
        htf.text = slide_data["highlight"]
        htf.paragraphs[0].font.size = Pt(14)
        htf.paragraphs[0].font.bold = True
        htf.paragraphs[0].font.color.rgb = WHITE
        htf.paragraphs[0].font.name = "Inter"
        htf.vertical_anchor = MSO_ANCHOR.MIDDLE

    add_footer(slide, theme, f"KM0 Digital · Slide {slide_num}", lang)


def build_two_col_slide(
    prs: Presentation, slide_data: dict, theme: Theme, lang: str, slide_num: int
) -> None:
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    set_slide_bg(slide, theme.bg)
    left = content_left(theme)

    if theme.slide_style == "light":
        add_accent_bar(slide, theme)
    elif theme.slide_style == "split":
        add_split_sidebar(slide, theme)

    title_box = slide.shapes.add_textbox(Inches(left), Inches(0.35), Inches(11.5 - left), Inches(0.8))
    tf = title_box.text_frame
    tf.text = slide_data["title"]
    style_title(tf, theme, size=28)

    col_w = (11.0 - left) / 2 - 0.15
    for idx, (title, items, accent) in enumerate(
        [
            (slide_data["left_title"], slide_data["left"], LIGHT_GRAY),
            (slide_data["right_title"], slide_data["right"], theme.accent),
        ]
    ):
        x = left + idx * (col_w + 0.3)
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(x),
            Inches(1.4),
            Inches(col_w),
            Inches(5.2),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = SURFACE if theme.slide_style != "dark" else RGBColor(0x15, 0x1D, 0x2E)
        card.line.color.rgb = accent

        hdr = slide.shapes.add_textbox(Inches(x + 0.2), Inches(1.55), Inches(col_w - 0.4), Inches(0.5))
        htf = hdr.text_frame
        htf.text = title
        htf.paragraphs[0].font.bold = True
        htf.paragraphs[0].font.size = Pt(18)
        htf.paragraphs[0].font.color.rgb = accent if idx == 1 else theme.title_color

        body = slide.shapes.add_textbox(Inches(x + 0.2), Inches(2.1), Inches(col_w - 0.4), Inches(4.2))
        btf = body.text_frame
        btf.word_wrap = True
        add_bullets(btf, items, theme, size=14)

    add_footer(slide, theme, f"KM0 Digital · Slide {slide_num}", lang)


def build_chips_slide(
    prs: Presentation, slide_data: dict, theme: Theme, lang: str, slide_num: int
) -> None:
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    set_slide_bg(slide, theme.bg)
    left = content_left(theme)

    if theme.slide_style == "light":
        add_accent_bar(slide, theme)
    elif theme.slide_style == "split":
        add_split_sidebar(slide, theme)

    title_box = slide.shapes.add_textbox(Inches(left), Inches(0.35), Inches(11.5 - left), Inches(0.8))
    tf = title_box.text_frame
    tf.text = slide_data["title"]
    style_title(tf, theme, size=30)

    x = left
    for i, chip in enumerate(slide_data["chips"]):
        w = 2.4
        chip_shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(x + i * 2.55),
            Inches(1.6),
            Inches(w),
            Inches(0.65),
        )
        colors = [ORANGE, MAGENTA, PURPLE, BLUE]
        chip_shape.fill.solid()
        chip_shape.fill.fore_color.rgb = colors[i % len(colors)]
        chip_shape.line.fill.background()
        ctf = chip_shape.text_frame
        ctf.text = chip
        ctf.paragraphs[0].font.size = Pt(16)
        ctf.paragraphs[0].font.bold = True
        ctf.paragraphs[0].font.color.rgb = WHITE
        ctf.paragraphs[0].alignment = PP_ALIGN.CENTER
        ctf.vertical_anchor = MSO_ANCHOR.MIDDLE

    body = slide.shapes.add_textbox(Inches(left), Inches(2.8), Inches(11.0 - left), Inches(3.5))
    btf = body.text_frame
    btf.text = slide_data["body"]
    btf.word_wrap = True
    style_body(btf, theme, size=20)

    add_footer(slide, theme, f"KM0 Digital · Slide {slide_num}", lang)


def build_timeline_slide(
    prs: Presentation, slide_data: dict, theme: Theme, lang: str, slide_num: int
) -> None:
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    set_slide_bg(slide, theme.bg)
    left = content_left(theme)

    if theme.slide_style == "light":
        add_accent_bar(slide, theme)
    elif theme.slide_style == "split":
        add_split_sidebar(slide, theme)

    title_box = slide.shapes.add_textbox(Inches(left), Inches(0.35), Inches(11.5 - left), Inches(0.8))
    tf = title_box.text_frame
    tf.text = slide_data["title"]
    style_title(tf, theme, size=28)

    y = 1.35
    for i, (label, desc) in enumerate(slide_data["items"]):
        dot = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            Inches(left),
            Inches(y + 0.05),
            Inches(0.18),
            Inches(0.18),
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = [ORANGE, MAGENTA, PURPLE, BLUE, ORANGE, MAGENTA][i]
        dot.line.fill.background()

        lbl = slide.shapes.add_textbox(Inches(left + 0.35), Inches(y - 0.05), Inches(1.5), Inches(0.35))
        ltf = lbl.text_frame
        ltf.text = label
        ltf.paragraphs[0].font.bold = True
        ltf.paragraphs[0].font.size = Pt(14)
        ltf.paragraphs[0].font.color.rgb = theme.accent

        desc_box = slide.shapes.add_textbox(Inches(left + 1.7), Inches(y - 0.05), Inches(9.0 - left), Inches(0.45))
        dtf = desc_box.text_frame
        dtf.text = desc
        dtf.paragraphs[0].font.size = Pt(14)
        dtf.paragraphs[0].font.color.rgb = theme.body_color
        y += 0.85

    add_footer(slide, theme, f"KM0 Digital · Slide {slide_num}", lang)


def build_cta_slide(
    prs: Presentation, slide_data: dict, theme: Theme, lang: str, slide_num: int
) -> None:
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    set_slide_bg(slide, NAVY if theme.slide_style != "light" else theme.bg)

    if theme.slide_style == "light":
        band = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
        )
        band.fill.solid()
        band.fill.fore_color.rgb = NAVY
        band.line.fill.background()

    add_logo(slide, 5.6, 0.4, 1.2)

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(0.8))
    tf = title_box.text_frame
    tf.text = slide_data["title"]
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.size = Pt(34)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.name = "Inter"

    body = slide.shapes.add_textbox(Inches(2.0), Inches(2.8), Inches(9.3), Inches(3.5))
    btf = body.text_frame
    btf.word_wrap = True
    for i, item in enumerate(slide_data["bullets"]):
        p = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
        p.text = f"→  {item}"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0xE5, 0xE7, 0xEB)
        p.font.name = "Inter"
        p.space_after = Pt(12)
        p.alignment = PP_ALIGN.LEFT

    add_footer(slide, theme, "KM0 Digital · AMVARA CONSULTING S.L.", lang)


def build_presentation(locale: str) -> Path:
    data = CONTENT[locale]
    theme = THEMES[locale]
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    builders = {
        "title": build_title_slide,
        "content": build_content_slide,
        "two_col": build_two_col_slide,
        "chips": build_chips_slide,
        "timeline": build_timeline_slide,
        "cta": build_cta_slide,
    }

    for i, slide_data in enumerate(data["slides"], start=1):
        slide_type = slide_data["type"]
        if slide_type == "title":
            builders["title"](prs, slide_data, theme, data["lang_label"])
        else:
            builders[slide_type](prs, slide_data, theme, data["lang_label"], i)

    out_pptx = OUT / f"{data['filename']}.pptx"
    prs.save(str(out_pptx))
    return out_pptx


def convert_to_pdf(pptx_path: Path) -> Path:
    pdf_path = pptx_path.with_suffix(".pdf")
    cmd = [
        "libreoffice",
        "--headless",
        "--convert-to",
        "pdf",
        "--outdir",
        str(OUT),
        str(pptx_path),
    ]
    subprocess.run(cmd, check=True, capture_output=True)
    return pdf_path


def main() -> None:
    generated: list[tuple[str, Path, Path]] = []
    for locale in ("ca", "es", "en"):
        pptx = build_presentation(locale)
        pdf = convert_to_pdf(pptx)
        generated.append((locale, pptx, pdf))
        print(f"✓ {locale.upper()}: {pptx.name} + {pdf.name}")

    print("\nAll files in:", OUT)
    for locale, pptx, pdf in generated:
        theme = THEMES[locale]
        print(f"  [{locale}] Theme '{theme.name}':")
        print(f"       PPT: {pptx}")
        print(f"       PDF: {pdf}")


if __name__ == "__main__":
    main()
